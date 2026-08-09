from __future__ import annotations

import csv
import json
from pathlib import Path
from typing import Any, Dict, List


def extract_txt(path: str | Path) -> Dict[str, Any]:
    p = Path(path)
    text = p.read_text(encoding="utf-8", errors="replace")
    return {"text": text, "tables": [], "metadata": {"suffix": p.suffix.lower()}}


def extract_csv(path: str | Path) -> Dict[str, Any]:
    p = Path(path)
    rows: List[List[str]] = []
    with p.open("r", encoding="utf-8-sig", errors="replace", newline="") as f:
        for row in csv.reader(f):
            rows.append([str(x) for x in row])
    text = "\n".join(" | ".join(r) for r in rows)
    return {"text": text, "tables": [rows], "metadata": {"suffix": p.suffix.lower()}}


def extract_json(path: str | Path) -> Dict[str, Any]:
    p = Path(path)
    data = json.loads(p.read_text(encoding="utf-8", errors="replace"))
    text = json.dumps(data, ensure_ascii=False, indent=2)
    return {"text": text, "tables": [], "metadata": {"suffix": p.suffix.lower()}}


def extract_pdf(path: str | Path) -> Dict[str, Any]:
    import fitz  # PyMuPDF

    p = Path(path)
    doc = fitz.open(p)
    pages = []
    for i, page in enumerate(doc):
        pages.append({"page": i + 1, "text": page.get_text("text")})
    text = "\n\n".join(x["text"] for x in pages)
    return {"text": text, "pages": pages, "tables": [], "metadata": {"pages": len(pages)}}


def extract_docx(path: str | Path) -> Dict[str, Any]:
    from docx import Document

    p = Path(path)
    doc = Document(p)
    paragraphs = [x.text for x in doc.paragraphs]
    tables = []
    for table in doc.tables:
        rows = [[cell.text for cell in row.cells] for row in table.rows]
        tables.append(rows)
    text = "\n".join(paragraphs)
    return {"text": text, "tables": tables, "metadata": {"paragraphs": len(paragraphs)}}


def extract_xlsx(path: str | Path) -> Dict[str, Any]:
    from openpyxl import load_workbook

    p = Path(path)
    wb = load_workbook(p, data_only=True, read_only=True)
    tables = []
    text_parts = []
    for ws in wb.worksheets:
        rows = []
        for row in ws.iter_rows(values_only=True):
            values = ["" if v is None else str(v) for v in row]
            rows.append(values)
        tables.append({"sheet": ws.title, "rows": rows})
        text_parts.append(f"# {ws.title}\n" + "\n".join(" | ".join(r) for r in rows))
    return {"text": "\n\n".join(text_parts), "tables": tables, "metadata": {"sheets": wb.sheetnames}}


def extract_pptx(path: str | Path) -> Dict[str, Any]:
    from pptx import Presentation

    p = Path(path)
    prs = Presentation(p)
    slides = []
    for idx, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                texts.append(shape.text)
        slides.append({"slide": idx, "text": "\n".join(texts)})
    text = "\n\n".join(f"## Slide {s['slide']}\n{s['text']}" for s in slides)
    return {"text": text, "pages": slides, "tables": [], "metadata": {"slides": len(slides)}}


def extract_document(path: str | Path) -> Dict[str, Any]:
    p = Path(path)
    suffix = p.suffix.lower()
    if suffix == ".pdf":
        return extract_pdf(p)
    if suffix == ".docx":
        return extract_docx(p)
    if suffix == ".xlsx":
        return extract_xlsx(p)
    if suffix == ".pptx":
        return extract_pptx(p)
    if suffix in {".txt", ".md"}:
        return extract_txt(p)
    if suffix == ".csv":
        return extract_csv(p)
    if suffix == ".json":
        return extract_json(p)
    raise ValueError(f"Unsupported document format: {suffix}")
